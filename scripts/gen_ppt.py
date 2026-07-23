#!/usr/bin/env python3
"""Generate defense PPT from v10.0 outline."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD   = RGBColor(0x25, 0x25, 0x3A)
ACCENT    = RGBColor(0x6C, 0x9C, 0xFF)
GREEN     = RGBColor(0x4E, 0xC9, 0xB0)
RED       = RGBColor(0xF8, 0x71, 0x71)
YELLOW    = RGBColor(0xFF, 0xD9, 0x3D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xA0, 0xA0, 0xB0)
DIM       = RGBColor(0x60, 0x60, 0x75)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, x=0.8, y=0.3, w=11.7, h=0.8, size=36):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    return tf

def add_subtitle(slide, text, x=0.8, y=1.1, w=11.7, h=0.5, size=18, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf

def add_body(slide, text, x=0.8, y=1.7, w=11.7, h=5.0, size=16, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(4)
    return tf

def add_card(slide, text, x, y, w, h, bg=BG_CARD, size=14, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf

def add_table(slide, headers, rows, x, y, w, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(0.4 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
    return table_shape

# ═══════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_title(slide, "F2C 全覆盖路径规划系统", y=1.8, size=44)
add_subtitle(slide, "从 Demo 82% 到 v10.0 十全十美", y=2.7, size=24, color=ACCENT)
add_subtitle(slide, "17 天 · 145 次迭代 · 21/21 GATE PASS · hole_crossings = 0", y=3.3, size=18)
add_subtitle(slide, "2026-07-23", y=4.2, size=16, color=DIM)

# ═══════════════════════════════════════════
# SLIDE 2: Problem
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "起点：Demo 基线 ~82%，五大核心挑战")
add_body(slide, """Demo (07/06): F2C 原生 RoutePlannerBase + OR-Tools AUTOMATIC TSP，覆盖率 ~82%，结果不可复现""", size=14, color=GRAY)
add_table(slide,
    ["挑战", "表现", "严重度"],
    [["孔洞穿越", "路径直接穿过障碍物 (S3/S6/S7)", "P0"],
     ["Cell 碎片化", "sweep 分解过度切分", "P0"],
     ["边界漏扫", "S1 顶部覆盖率 97%", "P1"],
     ["TSP 随机性", "AUTOMATIC 每次结果不同, N3 44-51 分波动", "P1"],
     ["斜边三角死区", "矩形 swath 无法贴合斜边", "P1"]],
    x=0.8, y=2.2, w=11.7,
    col_widths=[2.0, 5.7, 4.0])

# ═══════════════════════════════════════════
# SLIDE 3: Product constraints
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "产品约束：设计决策的边界条件")
add_subtitle(slide, "来源: docs/产品需求待办.md — 决定了什么能做、什么不能做")
add_table(slide,
    ["约束", "确认内容", "设计影响"],
    [["应用场景", "室内扫地机器人，非农机", "硬边界墙壁，覆盖优先于效率"],
     ["平台", "C++ / 嵌入式", "算法不能太重，不能依赖 GPU"],
     ["驱动方式", "差速驱动，可原地转", "turn_planner_type=direct (零半径)"],
     ["边界类型", "绝大部分 closed 硬边界", "swath 端点内缩留安全距离"],
     ["覆盖目标", "≥ 99.5% (硬门槛)", "覆盖不达标 = 方案无效"],
     ["Swath 重叠", "3% 故意重叠", "swath_overlap_ratio=0.03，防定位误差漏缝"],
     ["有效间距", "0.90 × 0.97 = 0.873m", "评估时扣除计划重叠，不扣分"]],
    x=0.8, y=1.6, w=11.7,
    col_widths=[2.2, 4.5, 5.0])

# ═══════════════════════════════════════════
# SLIDE 4: Timeline
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "演进总览：五阶段 · 145 commits · 7→21 场景")

phases = [
    ("基础期\n07/08-12", "V0 Demo\n~82%", "候选headland\n角度优化\nSweep分解"),
    ("架构期\n07/13-15", "v8 稳定版\n98.6%", "评分重构\ngreedy sort\nSnake优化"),
    ("质量期\n07/16-17", "v9.6 门控\n24 tests", "实车参数\n全项目审查\n边界填充重构"),
    ("硬化期\n07/18-20", "v9.11 穿洞\n硬约束 1e9", "C-space退役\n空隙分类器\nPhase 4A"),
    ("收敛期\n07/21-23", "v10.0\n21/21 PASS", "PATH_CHEAPEST_ARC\n微型去冗余\nScheme 2"),
]
for i, (title, milestone, details) in enumerate(phases):
    x = 0.5 + i * 2.5
    add_card(slide, title, x, 1.5, 2.3, 1.2, size=14)
    add_card(slide, milestone, x, 2.9, 2.3, 1.2, bg=ACCENT, size=14)
    add_card(slide, details, x, 4.3, 2.3, 2.5, size=12, color=GRAY)

# ═══════════════════════════════════════════
# SLIDE 5: Story 1 - greedyCellOrder
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事一：从固定极角序到贪心出口驱动")
add_body(slide, """\
为什么改: 旧方案用固定极角序（绕孔洞质心角度排序 cell），不考虑实际端点位置 → 跨场景长距离连接
\
greedyCellOrder 核心逻辑:
  1. 从 C0 出口开始，贪心选最近未访问 cell 端点
  2. 同时决定: 下一个 cell + 入口方向 + 内部覆盖方向 + 新出口
  3. 4 entry × 4 exit = 16 候选 → DP 选最优链
  4. 有孔洞: 圆形绕洞序 + 贪心方向反转；无孔洞: 纯贪心
  5. 穿洞惩罚: +1000（后强化为 +1e9）
\
效果: notched 由 3 条 27m 穿洞对角线 → 0
取舍: 保留极角序作为孔洞场景的 DP 遍历序基础，混合策略一直保留到 v10.0""", size=15)

# ═══════════════════════════════════════════
# SLIDE 6: Story 2 - Sweep + compensation chain
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事二：Sweep 分解 — 最大的架构决定")
add_body(slide, """\
为什么: F2C 原生 Boustrophedon 面对孔洞产生 O(n²) 碎片
结果: S3 42→6 cells (-86%), S6 +3.2% 覆盖率

但产生了连锁反应 — 整个项目最深层的架构认知:
  Sweep 分解 (替代 F2C Boustrophedon)
    → 对多边形朝向敏感          → Phase 4A 倾斜 sweep 补偿
    → 产生碎片化条带             → Cell 合并 (两轮安全合并) 补偿
    → 斜边三角间隙               → 斜边检测 + 独立角度补偿
    → 全局角度对部分 cell 不优   → Per-cell veto 补偿
    → Veto 改变方向后可能穿洞    → 1e9 硬约束补偿

教训: Sweep 打破了 F2C 原生假设；每个断裂点需要一个补偿层；补偿层互相耦合
取舍: O(n²)→O(n) 的收益值得后续 5 层补偿""", size=14)

# ═══════════════════════════════════════════
# SLIDE 7: Story 3 - Hole crossing battle
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事三：穿洞之战 — 三个回合，从软到硬")
add_table(slide,
    ["回合", "方案", "结果", "教训"],
    [["第一回合\n07/13", "后置打补丁: 检测穿洞段 + 绕行点", "方案 B segfault\n回退", "后置修复药不对症"],
     ["第二回合\n07/14", "Cell 连接反转 + 穿洞检测", "notched +26 分\n但仍有场景穿", "部分解决"],
     ["第三回合\n07/20", "1e9 硬约束: +1000 → +1e9", "8/8 hole_crossings=0\n零退化", "最小修复原则:\n只改惩罚"]],
    x=0.8, y=1.6, w=11.7,
    col_widths=[1.8, 3.5, 3.0, 3.4])
add_body(slide, """\
核心教训: 组合优化中软约束在退化情况下失效 — 正常时 +1000 够用，全部候选穿洞时 DP 没得选""",
    y=4.5, size=14, color=YELLOW)

# ═══════════════════════════════════════════
# SLIDE 8: Story 4 - Cell merge
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事四：Cell 合并的诱惑与危险")
add_table(slide,
    ["方案", "做法", "结果", "为什么失败"],
    [["共享边检测", "只合并共享边 cell", "零合并", "GDAL 后几乎无共享边"],
     ["激进两轮合并", "Pass2 无孔洞保护", "N3 98.1%", "cell 跨洞合并"],
     ["最终方案 ✅", "v9.11 + x-span strip", "20/21 PASS", "interior ring 把关"]],
    x=0.8, y=1.6, w=11.7,
    col_widths=[2.2, 3.0, 2.5, 4.0])
add_body(slide, """\
效果: S7 20→9 cells, N2 11→5 cells, S3 6→3 cells
但 N3 仍然碎片化 (98.1%) — 引出下一个故事""", y=4.3, size=14, color=YELLOW)

# ═══════════════════════════════════════════
# SLIDE 9: Story 5 - X-cut rollbacks
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事五：X-cut 的三次回滚")
add_table(slide,
    ["迭代", "方案", "N3", "退化"],
    [["1: 全局 X-cut", "所有孔洞 X 边界加入分割线", "100% ✅", "❌ S7/N13/S3"],
     ["2: 条件触发", "Y重叠+X不重叠孔洞对才切", "100% ✅", "❌ S7/N13"],
     ["3: 更严条件", "≥2×2 或 1×3 才切", "—", "❌ 老板:面向测试场景编程"]],
    x=0.8, y=1.6, w=11.7,
    col_widths=[2.0, 4.0, 2.0, 3.7])
add_body(slide, """\
全部回滚。正确做法 (Scheme 2): 只在孔洞条带内部 X-cut → N3 98.1%→99.97%, +15.2 分
教训: 让方案驱动，不要自行加东西；面向测试场景编程是坏味道""", y=4.3, size=14, color=YELLOW)

# ═══════════════════════════════════════════
# SLIDE 10: Story 6 - Boundary fill
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "故事六：边界填充的 6 轮进化")
add_body(slide, """\
第一版: 沿外环+孔洞边界生成填充 swath
问题1: 去重大激进 (0.3*cov_width) → S4 75%           修复1: 移除去重 → S4 75%→97.4% (+22%)
问题2: touches_outer 全局门控误杀贴边 cell              修复2: 改为逐边独立判断
问题3: 边界填充被 genRoute 踢到 Route 尾部 → 越界连接  修复3: 填充进入 genRoute 前 → notched +60.8 分
问题4: 相邻 cell 各自生成同一缝隙填充 → S3 重叠翻倍     修复4: pruneRedundantCellSeamFills
问题5: C-space 孔洞边界重复填充                         修复5: isCoveredByExistingSwath
\
教训: "贴边补一刀"看似简单，6 轮迭代才稳定。每一轮修一个问题，同时引入一个新问题。边界处理没有银弹。""", size=14)

# ═══════════════════════════════════════════
# SLIDE 11: Final results
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "最终成果：21/21 GATE PASS")
add_table(slide,
    ["指标", "数值"],
    [["场景数", "21 (含 13 个新增极端场景)"],
     ["有效覆盖率", "99.94% ~ 100%"],
     ["GATE PASS", "21/21 ✅"],
     ["hole_crossings", "0 (全部场景)"],
     ["S7 工厂车间", "0% → 100%"],
     ["N11_whole", "59% → 100%"],
     ["平均得分", "~75"],
     ["编译", "0 error 0 warning"],
     ["单元测试", "47/47 pass"]],
    x=0.8, y=1.6, w=7.0,
    col_widths=[3.0, 4.0])

# ═══════════════════════════════════════════
# SLIDE 12: Methodology - Evaluation loop
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "怎么做 ①：测评驱动的优化闭环")
add_body(slide, """\
评分公式:
  CoverageGate = clamp((c - 0.90) / 0.095, 0, 1)³     ← 90%=0, 99.5%=满分
  Score = CoverageGate × EfficiencyScore               ← 工作比+距离+转弯+重叠+耗时

设计意图: 覆盖率不达标 → 一票否决 (gate=0)；99.5% 以上满分 → 防止钻 99% 空子；效率分在覆盖合格后才参与比较

自动化工具链:
  quick_dev_test.sh (5 场景) → 日常迭代秒级反馈
  batch_test_v2.sh (21 场景) → 提交前 15 分钟全量验证
  n3_repeat_test.sh (单场景) → TSP 方差验证

案例: N11 59% 回归 → Codex 静审 + Claude 代码审查 → 1h 定位 → 修复
没有这套体系，N11 的 59% 可能几周后才被发现，届时已经不知道哪个改动引入的""", size=14)

# ═══════════════════════════════════════════
# SLIDE 13: Methodology - Modeling
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "怎么做 ②：工程建模的务实取舍")
add_body(slide, """\
三层模型，精度逐级递减:
  物理足迹层 (碰撞检测)  ← 0.05m 采样，逐姿态矩形+双轮检查，仅后处理
  覆盖模型层 (swath生成) ← coverage_width=0.90，清洁带=完美矩形，清洁长度=车长
  拓扑模型层 (分解/合并) ← robot_width=0.75，圆形包络近似，故意保守

已知简化 & 为什么可以接受:
  • 清洁长度=车长 → swath 起点存在未清洁死角 → 仅在端点，面积 <0.1%
  • 转弯覆盖未建模 → 转弯扫过面积不计入 → 保守方向，实际覆盖 ≥ 报告值
  • 转弯碰撞已检测 → 安全不妥协，覆盖可保守
  • 3% Swath 重叠 (swath_overlap_ratio=0.03) → 有效间距 0.873m → 防定位误差漏缝

碰撞安全与拓扑规划参数分离:
  robot_width 早期一人分饰六角 → 引入 PhysicalFootprintParams 独立负责碰撞 → 解耦实验证实隐藏依赖 → 改从另一角度解决""", size=14)

# ═══════════════════════════════════════════
# SLIDE 14: Methodology - Vibe Coding
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "怎么做 ③：Vibe Coding 的工作流纪律")
add_body(slide, """\
纪律一: 模块化解耦
  polygon_planner_node (ROS2) → PlannerCore::plan() (纯 C++)
    ├ decomposer (~400行)  ├ swath_generator (~600行)  ├ boundary_filler (~500行)
    ├ path_planner (~600行) └ coverage_evaluator (~700行)
  每个模块 <800 行，改算法不动 ROS → AI 友好

纪律二: 快速验证通道
  F2C Python 接口 → 不编译也能验证 (秒级 vs 分钟级)
  quick_dev_test.sh (5场景) → 日常迭代  /  batch_test_v2.sh (21场景) → 提交前

纪律三: AI+人都能读的数据界面
  JSON 输出 (坐标/swath/连接) → AI 精准定位
  42 张可视化 PNG (覆盖+Cell) → 人一眼判断
  连接 JSON 从仅存 from_cell/to_cell → 存 route_waypoints (1087条全部含实际路径)

纪律四: 终端集中操作 → SSH+SCP，所有操作可追溯、可复现""", size=14)

# ═══════════════════════════════════════════
# SLIDE 15: Methodology - Governance
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "怎么做 ④：并行探索的工程治理")
add_table(slide,
    ["分支", "用途", "状态"],
    [["master", "稳定主线 (PlannerCore only)", "✅ 当前"],
     ["legacy-dual-pipeline", "旧双管线存档", "GitHub 备份"],
     ["main-stable", "v9.11 + 可见图路由", "穿洞问题，参考"],
     ["feat/param-decouple-0721", "参数解耦实验", "实验存档"],
     ["codex/experiment-...", "Codex 安全连接捷径", "实验"]],
    x=0.8, y=1.5, w=7.0, col_widths=[3.0, 2.5, 1.5])
add_body(slide, """\
多分支并行 → 不同风险等级同时推进，互不干扰
tag 基线 (v9.12-stable, v10.0) → 任何时刻可回到已知状态
GitHub 远程备份 → 不怕 VM 炸、不怕本地改坏
145 commits 可追溯 → "谁改了什么为什么"全有记录
→ AI 辅助开发的安全网：没有分支隔离和版本管理，vibe coding 就是灾难""", y=4.4, size=14)

# ═══════════════════════════════════════════
# SLIDE 16: Methodology - Human-AI
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "怎么做 ⑤：人机协作 + 工作流")
add_body(slide, """\
分工: 老板(架构/参数/场景) · Claude(分析/审查/对比) · Codex(实现/可视化/脚本)

交叉审查: 声明→证据→定量推演  |  两个 AI 独立审查，不互相背书
Codex 写代码 → Claude 验证 → 发现问题 → Codex 修复 → Claude 再验证

工作流:
  Windows (Claude+Codex+Git) ←ssh→ VM (ROS2+colcon+batch_test)
  改代码 → sync_and_build.sh → 编译 → batch_test_v2.sh → 拉回 21 PNG + JSON
  → 交叉审查 → 决策: 接受/回滚/继续

核心原则: 重复性工作给工具，判断力留给自己。AI 是工具不是答案。""", size=14)

# ═══════════════════════════════════════════
# SLIDE 17: Future
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "改进方向")
add_table(slide,
    ["优先级", "问题", "方案"],
    [["P1", "N3 重叠 101% (21 cells 碎片化)", "有条件 Cell 合并 (路由层孔洞保护)"],
     ["P1", "S7 重叠 53.8% (128 turns, 896.9m)", "安全检查后直连替代微 U 绕行"],
     ["P1", "N9 重叠 46.4% (货架边缘重复补线)", "补线去重"],
     ["P1", "Swath 层共线合并", "同方向+同自由空间带 swath 归一化"],
     ["P2", "Cell-block 贪心架构", "方向正确，需 visibility graph 安全连接"],
     ["P2", "线程安全 / 测试覆盖 / WSL2 迁移", "当前零 mutex，单线程 safe"],
     ["—", "确认舍弃", "A*路由 / Snake排序 / X-cut全局版 / 参数解耦"]],
    x=0.8, y=1.5, w=11.7,
    col_widths=[0.8, 4.5, 6.4])

# ═══════════════════════════════════════════
# SLIDE 18: Summary
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "总结", size=40)
add_body(slide, """\
成  果: 82% → ~100% 有效覆盖率 · 21/21 GATE PASS · hole_crossings = 0
         S7 工厂车间 0%→100% · N11 59%→100%

方  法: 测评驱动 · 模块化 · 交叉审查 · 工程治理
         145 commits · 多分支并行 · tag 基线 · 数据 > 直觉

理  念: 回滚 ≠ 失败 · 局部优化不能以全局稳定性为代价
         AI 是工具不是答案 · 重复性工作给工具，判断力留给自己""", y=2.0, size=18, bold=False)

# ── Save ──
output_path = "docs/v10.0-defense-ppt.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
